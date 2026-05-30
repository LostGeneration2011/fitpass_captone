from __future__ import annotations

from pathlib import Path
import re

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
from pptx.enum.text import PP_ALIGN


INPUT_MD = Path(r"c:\New folder\fitpass\FITPASS_PRESENTATION_CONTENT.md")
OUTPUT_PPTX = Path(r"c:\New folder\fitpass\FITPASS_PRESENTATION_OUTPUT.pptx")


SLIDE_HEADER_RE = re.compile(r"^##\s+SLIDE\s+\d+\s+[—-]\s+(.+)$", re.IGNORECASE)

ACCENT = RGBColor(13, 71, 161)
ACCENT_DARK = RGBColor(8, 43, 102)
TEXT_PRIMARY = RGBColor(33, 37, 41)
TEXT_MUTED = RGBColor(73, 80, 87)
LIGHT_BG = RGBColor(245, 248, 252)


def clean_markdown_line(line: str) -> str:
    line = line.rstrip()
    if not line:
        return ""

    # Convert markdown table rows into concise plain text.
    if line.strip().startswith("|") and line.strip().endswith("|"):
        cells = [c.strip() for c in line.strip().strip("|").split("|")]
        if cells and all(set(c) <= {"-", ":"} for c in cells):
            return ""
        if len(cells) >= 3:
            return f"{cells[0]} | {cells[1]}: {cells[2]}"
        if len(cells) == 2:
            return f"{cells[0]}: {cells[1]}"
        if len(cells) == 1:
            return cells[0]

    # Remove common Markdown markers while keeping readable text.
    line = re.sub(r"^[-*]\s+", "", line)
    line = re.sub(r"^\d+\.\s+", "", line)
    line = re.sub(r"^#{1,6}\s+", "", line)
    line = line.replace("**", "")
    line = line.replace("`", "")
    line = re.sub(r"\[(.*?)\]\(.*?\)", r"\1", line)
    return line


def parse_slides(markdown_text: str) -> list[tuple[str, list[str]]]:
    slides: list[tuple[str, list[str]]] = []
    current_title: str | None = None
    current_lines: list[str] = []

    in_code_block = False

    for raw_line in markdown_text.splitlines():
        line = raw_line.rstrip("\n")

        if line.strip().startswith("```"):
            in_code_block = not in_code_block
            continue

        header_match = SLIDE_HEADER_RE.match(line.strip())
        if header_match and not in_code_block:
            if current_title is not None:
                slides.append((current_title, current_lines))
            current_title = header_match.group(1).strip()
            current_lines = []
            continue

        if current_title is None:
            continue

        if line.strip() == "---":
            continue

        cleaned = clean_markdown_line(line)
        if cleaned:
            current_lines.append(cleaned)

    if current_title is not None:
        slides.append((current_title, current_lines))

    return slides


def add_theme_background(slide) -> None:
    background = slide.background
    background.fill.solid()
    background.fill.fore_color.rgb = LIGHT_BG

    top_bar = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.RECTANGLE,
        Inches(0),
        Inches(0),
        Inches(13.333),
        Inches(0.45),
    )
    top_bar.fill.solid()
    top_bar.fill.fore_color.rgb = ACCENT
    top_bar.line.fill.background()


def add_footer(slide, index: int) -> None:
    footer = slide.shapes.add_textbox(Inches(0.5), Inches(7.1), Inches(12.4), Inches(0.25))
    tf = footer.text_frame
    p = tf.paragraphs[0]
    p.text = f"FitPass Presentation | Slide {index}"
    p.font.size = Pt(10)
    p.font.color.rgb = TEXT_MUTED
    p.alignment = PP_ALIGN.RIGHT


def chunk_lines(lines: list[str], max_lines: int = 10) -> list[list[str]]:
    chunks: list[list[str]] = []
    current: list[str] = []

    for line in lines:
        if not line:
            continue
        if len(current) >= max_lines:
            chunks.append(current)
            current = []
        current.append(line)

    if current:
        chunks.append(current)

    return chunks or [["(No content provided)"]]


def add_cover_slide(prs: Presentation, title: str, lines: list[str], index: int) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    add_theme_background(slide)

    title_shape = slide.shapes.title
    subtitle = slide.placeholders[1]

    title_shape.text = "FITPASS"
    title_para = title_shape.text_frame.paragraphs[0]
    title_para.font.size = Pt(48)
    title_para.font.bold = True
    title_para.font.color.rgb = ACCENT_DARK
    title_para.alignment = PP_ALIGN.CENTER

    subtitle.text = "\n".join(lines[:6]) if lines else title
    for p in subtitle.text_frame.paragraphs:
        p.font.size = Pt(18)
        p.font.color.rgb = TEXT_PRIMARY
        p.alignment = PP_ALIGN.CENTER

    badge = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE,
        Inches(5.2),
        Inches(0.7),
        Inches(2.9),
        Inches(0.45),
    )
    badge.fill.solid()
    badge.fill.fore_color.rgb = ACCENT
    badge.line.fill.background()
    badge_text = badge.text_frame.paragraphs[0]
    badge_text.text = "Smart Gym Management System"
    badge_text.font.size = Pt(12)
    badge_text.font.bold = True
    badge_text.font.color.rgb = RGBColor(255, 255, 255)
    badge_text.alignment = PP_ALIGN.CENTER

    add_footer(slide, index)


def add_content_slide(prs: Presentation, title: str, lines: list[str], index: int) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    add_theme_background(slide)
    slide.shapes.title.text = title

    title_tf = slide.shapes.title.text_frame
    title_p = title_tf.paragraphs[0]
    title_p.font.bold = True
    title_p.font.size = Pt(31)
    title_p.font.color.rgb = ACCENT_DARK

    content = slide.placeholders[1].text_frame
    content.clear()

    for idx, line in enumerate(lines):
        p = content.paragraphs[0] if idx == 0 else content.add_paragraph()
        p.text = line
        p.level = 0
        p.font.size = Pt(18 if idx == 0 else 16)
        p.font.color.rgb = TEXT_PRIMARY
        if line.endswith(":"):
            p.font.bold = True

    add_footer(slide, index)


def build_presentation() -> Path:
    if not INPUT_MD.exists():
        raise FileNotFoundError(f"Missing input markdown: {INPUT_MD}")

    markdown_text = INPUT_MD.read_text(encoding="utf-8")
    slides = parse_slides(markdown_text)
    if not slides:
        raise ValueError("No slide sections found in markdown.")

    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    rendered_count = 0
    for idx, (title, lines) in enumerate(slides):
        if idx == 0:
            rendered_count += 1
            add_cover_slide(prs, title, lines, rendered_count)
            continue

        for chunk_idx, chunk in enumerate(chunk_lines(lines, max_lines=10)):
            rendered_count += 1
            chunk_title = title if chunk_idx == 0 else f"{title} (cont.)"
            add_content_slide(prs, chunk_title, chunk, rendered_count)

    prs.save(str(OUTPUT_PPTX))
    return OUTPUT_PPTX


if __name__ == "__main__":
    out = build_presentation()
    print(f"Created: {out}")
